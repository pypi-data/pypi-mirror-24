import io
import os
import sys
import copy
import json
import collections
import pandas as pd
from . import commands
from pptx import Presentation
from orderedattrdict import AttrDict
from pptx.shapes.shapetree import SlideShapes
from orderedattrdict.yamlutils import AttrDictYAMLLoader
from .utils import merge, parse_command_line, load_yaml, load_data, is_slide_allowed, is_group

_folder = os.path.dirname(os.path.abspath(__file__))
with io.open(os.path.join(_folder, 'release.json'), encoding='utf-8') as _release_file:
    _release = json.load(_release_file)
    __version__ = _release['version']


def commandline():
    '''
    Runs PPTGen from the command line.
    This is called via setup.py console_scripts.
    Though a trivial function, is is kept different from run_commands to allow
    unit testing of run_commands.
    '''
    run_commands(sys.argv[1:], pptgen)


def run_commands(commands, callback):
    '''
    For example::

        run_commands(['a.yaml', 'b.yaml', '--x=1'], method)

    will do the following:

    - Load a.yaml into config
        - Set config['a'] = 1
        - Change to directory where a.yaml is
        - Call method(config)
    - Load b.yaml into config
        - Set config['a'] = 1
        - Change to directory where b.yaml is
        - Call method(config)

    Command line arguments are passed as ``commands``.
    Callback is a function that is called for each config file.
    '''
    args = parse_command_line(commands)
    original_path = os.getcwd()
    for config_file in args.pop('_'):
        config = load_yaml(config_file, Loader=AttrDictYAMLLoader)
        config = merge(old=config, new=args, mode='overwrite')
        os.chdir(os.path.dirname(os.path.abspath(config_file)))
        try:
            callback(**config)
        finally:
            os.chdir(original_path)


def manage_data(data):
    '''
    Data management function.
    '''
    if isinstance(data, (pd.DataFrame, str,)):
        return json.loads(data.to_json(orient='records'))
    return data


def replicate_slides(data, prs, change, slide, slides_to_remove, index):
    '''
    Function to replicate slides.
    '''
    data = manage_data(data)
    copy_slide = copy.deepcopy(slide)
    slides_to_remove.append(index)
    # Stacking shapes if required.
    commands.stack_shapes(copy_slide.shapes, change, data)
    new_slide = commands.generate_slide(prs, copy_slide)
    change_shapes(copy_slide.shapes, change, data,
                  prs=prs, copy_slide=True,
                  source_slide=slide, new_slide=new_slide)
    copy_slide = None


def pptgen(source, target, **config):
    '''
    Process a configuration. This loads a Presentation from source, applies the
    (optional) configuration changes and saves it into target.
    '''
    data = AttrDict(load_data(config.pop('data', {})))
    prs = Presentation(source)
    slides = prs.slides
    # Loop through each change configuration
    slides_to_remove = []
    manage_slide_order = collections.defaultdict(list)

    for key, change in config.items():
        # Apply it to every slide
        slide_data = copy.deepcopy(data)
        if 'data' in change and change['data'] is not None:
            slide_data = eval('slide_data.{}'.format(change['data']))
        slide_data = manage_data(slide_data)
        for index, slide in enumerate(slides):
            # Restrict to specific slides, if specified
            if not is_slide_allowed(change, slide, index + 1):
                continue

            if 'replicate' in change and change['replicate'] is not None:
                is_grp = isinstance(slide_data, pd.core.groupby.DataFrameGroupBy)
                if isinstance(slide_data, collections.Iterable):
                    for _slide_data in slide_data:
                        _slide_data = _slide_data[1] if is_grp is True else _slide_data
                        replicate_slides(_slide_data, prs, change, slide, slides_to_remove, index)
                        # Creating dict mapping to order slides.
                        manage_slide_order[index + 1].append(len(prs.slides))
                else:
                    raise NotImplementedError()
            else:
                # Stacking shapes if required.
                commands.stack_shapes(slide.shapes, change, slide_data)
                change_shapes(slide.shapes, change, slide_data)

    indexes = []
    for key in sorted(manage_slide_order.keys()):
        indexes.append(manage_slide_order[key])

    matrix = list(map(list, zip(*indexes)))

    for indx_lst in matrix:
        for idx in indx_lst:
            src = prs.slides[idx - 1]
            slides_to_remove.append(idx - 1)
            copy_slide = copy.deepcopy(src)
            new_slide = commands.generate_slide(prs, copy_slide)
            dest = prs.slides.add_slide(new_slide)
            for shape in copy_slide.shapes:
                commands.copy_slide_elem(shape, dest)
            commands.add_new_slide(dest, src)

    removed_status = 0
    for sld_idx in set(slides_to_remove):
        commands.delete_slide(prs, (sld_idx - removed_status))
        for slide_num in manage_slide_order:
            manage_slide_order[slide_num] = [(i - 1) for i in manage_slide_order[slide_num]]
        removed_status += 1
    prs.save(target)


def change_shapes(collection, change, data,
                  prs=None, copy_slide=False,
                  source_slide=None, new_slide=None):
    '''
    Apply changes to a collection of shapes in the context of data.
    ``collection`` is a slide.shapes or group shapes.
    ``change`` is typically a dict of <shape-name>: commands.
    ``data`` is a dictionary passed to the template engine.
    '''
    dest = prs.slides.add_slide(new_slide) if copy_slide else None

    mapping = {}
    for shape in collection:
        if shape.name not in change:
            commands.copy_slide_elem(shape, dest)
            continue

        spec = change[shape.name]
        if shape.name not in mapping:
            mapping[shape.name] = 0

        shape_data = copy.deepcopy(eval(spec['data'])) if 'data' in spec else copy.deepcopy(data)
        if 'stack' in spec and spec.get('stack') is not None:
            shape_data = shape_data[mapping[shape.name]]
        mapping[shape.name] = mapping[shape.name] + 1

        # If the shape is a group, apply spec to each sub-shape
        if is_group(shape):
            sub_shapes = SlideShapes(shape.element, collection)
            change_shapes(sub_shapes, spec, shape_data)

        # Run commands in the spec
        for cmd, method in commands.cmdlist.items():
            if cmd in spec:
                method(shape, spec, shape_data)

        commands.copy_slide_elem(shape, dest)
    commands.add_new_slide(dest, source_slide)
