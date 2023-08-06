"""Python-PPTX customized module."""
import os
import six
import copy
import json
import logging
import requests
import platform
import datetime
import tempfile
import collections
import numpy as np
import pandas as pd
from . import stats
from . import fontwidth
from six import iteritems
from lxml import objectify
from pptx.util import Inches
from . import color as _color
from pptx.dml.color import RGBColor
from lxml.builder import ElementMaker
from tornado.template import Template
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import XyChartData
from pptx.chart.data import BubbleChartData
from six.moves.urllib_parse import urlparse

_template_cache = {}


def template(tmpl, data):
    """Execute tornado template."""
    if tmpl not in _template_cache:
        _template_cache[tmpl] = Template(tmpl, autoescape=None)
    # May not work with python 2.7
    return _template_cache[tmpl].generate(**data).decode()


def text(shape, spec, data):
    '''
    Replace entire text of shape with spec['text'].
    '''
    if not shape.has_text_frame:
        logging.error('"%s" is not a TextShape to apply text:', shape.name)
        return

    run_flag = True
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = template(spec['text'], data) if run_flag else ''
            run_flag = False


def replace(shape, spec, data):
    '''
    Replace keywords in shape using the dictionary at spec['replace'].
    '''
    if not shape.has_text_frame:
        logging.error('"%s" is not a TextShape to apply text:', shape.name)
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old, new in spec['replace'].items():
                run.text = run.text.replace(old, template(new, data))


def oval(shape, spec, data):
    """Function to animate oval shape type."""
    fill = shape.fill
    fill.solid()
    data['_color'] = _color
    bg_color = template(spec['oval']['background-color'], data)
    fill.fore_color.rgb = RGBColor.from_string(bg_color.replace('#', ''))


def image(shape, spec, data):
    '''
    Replace image with a different file specified in spec['image']
    '''
    image = template(spec['image'], data)
    # If it's a URL, use the requests library's raw stream as a file-like object
    if urlparse(image).netloc:
        r = requests.get(image)
        with tempfile.NamedTemporaryFile(delete=False) as handle:
            handle.write(r.content)
        new_img_part, new_rid = shape.part.get_or_add_image_part(handle.name)
        os.unlink(handle.name)
    else:
        new_img_part, new_rid = shape.part.get_or_add_image_part(image)
    old_rid = shape._pic.blip_rId
    shape._pic.blipFill.blip.rEmbed = new_rid
    shape.part.related_parts[old_rid].blob = new_img_part.blob


def stack_elements(replica, shape, stack=False, margin=None):
    '''
    Function to extend elements horizontally or vertically.
    '''
    config = {'vertical': {'axis': 'y', 'attr': 'height'},
              'horizontal': {'axis': 'x', 'attr': 'width'}}
    if stack is not None:
        grp_sp = shape.element
        # Adding a 15% margin between original and new object.
        margin = 0.50 if not margin else margin
        for index in range(replica - 1):
            # Adding a cloned object to shape
            extend_shape = copy.deepcopy(grp_sp)
            # Getting attributes and axis values from config based on stack.
            attr = config.get(stack, {}).get('attr', 0)
            axis = config.get(stack, {}).get('axis', 0)
            # Taking width or height based on stack value and setting a margin.
            metric_val = getattr(shape, attr)
            axis_val = getattr(extend_shape, axis)
            # Setting margin accordingly either vertically or horizontally.
            axis_pos = metric_val * (index + 1)
            set_attr = axis_val + axis_pos + (axis_pos * margin)
            # Setting graphic position of newly created object to slide.
            setattr(extend_shape, axis, int(set_attr))
            # Adding newly created object to slide.
            grp_sp.addnext(extend_shape)


def stack_shapes(collection, change, data):
    '''
    Function to stack Shapes if required.
    '''
    shape_data = data
    for shape in collection:
        if shape.name not in change:
            continue
        info = change[shape.name]
        if 'data' in info and info.get('stack') is not None:
            shape_data = eval(info.get('data'))
        stack_elements(len(shape_data), shape, stack=info.get('stack'),
                       margin=info.get('margin'))


def generate_slide(prs, source):
    '''
    Create a slide layout.
    '''
    layout_items_count = [
        len(layout.placeholders) for layout in prs.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return prs.slide_layouts[blank_layout_id]


def copy_slide_elem(shape, dest):
    '''
    Function to copy slide elements into a newly created slide.
    '''
    if dest:
        el = shape.element
        new_elem = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(new_elem, 'p:extLst')


def add_new_slide(dest, source_slide):
    '''
    Function to add a new slide to presentation.
    '''
    if dest:
        for key, value in six.iteritems(source_slide.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                dest.part.rels.add_relationship(value.reltype, value._target, value.rId)


def replicate_slide(change, prs, data):
    '''
    Add Duplicate slide to the presentation.
    '''
    if 'replicate-slides' in change and change['replicate-slides'] is not None:
        replica_conf = change['replicate-slides']['slide-number']
        for conf in replica_conf:
            for slide_number, replica_logic in conf.items():
                replica = template(replica_logic, data)
                for replica in range(int(replica) - 1):
                    source = prs.slides[slide_number - 1]
                    add_new_slide(prs, source)


def move_slide(presentation, old_index, new_index):
    '''
    Move a slide's index number.
    '''
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def delete_slide(presentation, index):
    '''
    Delete a slide from Presentation.
    '''
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def _prepare_data(config, data):
    '''
    Prepare chart data.
    '''
    if not isinstance(data[config['data']], (list, pd.DataFrame)):
        raise NotImplementedError()

    module_data = data[config['data']]

    if isinstance(data[config['data']], list):
        module_data = pd.DataFrame(data[config['data']])

    # Filtering data based on query
    if config.get('query') is not None and 'query' in config:
        module_data = module_data.query(config['query'])

    # Create a new dataset with the relevant columns
    if config.get('columns') is not None and 'columns' in config:
        result = pd.DataFrame()
        for name, column_info in config['columns'].items():
            if 'value' in column_info:
                result[name] = module_data[column_info['value']]
            elif 'formula' in column_info:
                result[name] = module_data.eval(column_info['formula'])
        return result

    return module_data


def _update_chart(info, data, chart_data, series_columns, chart='line'):
    '''
    Updating XyChart data.
    '''
    if chart == 'line':
        for series in series_columns:
            chart_data.add_series(
                series, tuple(data[series].dropna().tolist()))
        return chart_data

    for index, row in data.iterrows():
        for col in series_columns:
            series = chart_data.add_series(col)

            if chart == 'scatter':
                series.add_data_point(row[info['x']], row[col])

            elif chart == 'bubble':
                bubble_size = 1
                if info.get('size') is not None and 'size' in info:
                    bubble_size = row[info['size']]
                series.add_data_point(row[info['x']], row[col], bubble_size)

    return chart_data


def chart(shape, spec, data):
    '''
    Replacing chart Data.
    '''
    chart_type = None
    if hasattr(shape.chart, 'chart_type'):
        chart_type = '{}'.format(shape.chart.chart_type).split()[0]

    stacked_or_line = [
        'AREA', 'AREA_STACKED', 'AREA_STACKED_100', 'BAR_CLUSTERED',
        'BAR_OF_PIE', 'BAR_STACKED', 'BAR_STACKED_100', 'COLUMN_CLUSTERED',
        'COLUMN_STACKED', 'COLUMN_STACKED_100', 'CONE_BAR_CLUSTERED',
        'CONE_BAR_STACKED', 'CONE_BAR_STACKED_100', 'CONE_COL',
        'CONE_COL_CLUSTERED', 'CONE_COL_STACKED', 'CONE_COL_STACKED_100',
        'CYLINDER_BAR_CLUSTERED', 'CYLINDER_BAR_STACKED',
        'CYLINDER_BAR_STACKED_100', 'CYLINDER_COL', 'CYLINDER_COL_CLUSTERED',
        'CYLINDER_COL_STACKED', 'CYLINDER_COL_STACKED_100', 'LINE',
        'LINE_MARKERS', 'LINE_MARKERS_STACKED', 'LINE_MARKERS_STACKED_100',
        'LINE_STACKED', 'LINE_STACKED_100', 'PYRAMID_BAR_CLUSTERED',
        'PYRAMID_BAR_STACKED', 'PYRAMID_BAR_STACKED_100', 'PYRAMID_COL',
        'PYRAMID_COL_CLUSTERED', 'PYRAMID_COL_STACKED', 'RADAR_MARKERS',
        'PYRAMID_COL_STACKED_100', 'RADAR', 'RADAR_FILLED', 'PIE',
        'PIE_EXPLODED', 'PIE_OF_PIE']

    xy_charts = [
        'XY_SCATTER', 'XY_SCATTER_LINES', 'XY_SCATTER_LINES_NO_MARKERS',
        'XY_SCATTER_SMOOTH', 'XY_SCATTER_SMOOTH_NO_MARKERS']

    bubble_charts = ['BUBBLE', 'BUBBLE_THREE_D_EFFECT']

    if not chart_type:
        raise NotImplementedError()

    info = spec['chart']

    change_data = _prepare_data(info, data)
    series_cols = [x for x in change_data.columns if x != info['x']]
    # If chart type is stacked bar or line.
    if chart_type in stacked_or_line:
        # Initializing chart data
        chart_data = ChartData()
        chart_data.categories = change_data[info['x']].dropna().unique().tolist()
        change_data = _update_chart(info, change_data, chart_data, series_cols)

    # If chart type is scatter plot.
    elif chart_type in xy_charts:
        # Initializing chart data
        chart_data = XyChartData()
        change_data = _update_chart(info, change_data, chart_data, series_cols, chart='scatter')

    # If chart type is bubble chart.
    elif chart_type in bubble_charts:
        # Initializing chart data
        chart_data = BubbleChartData()
        change_data = _update_chart(info, change_data, chart_data, series_cols, chart='bubble')
    else:
        raise NotImplementedError()
    shape.chart.replace_data(chart_data)

    if spec['chart'].get('color'):
        colors = data['chart_data'][spec['chart']['color']].dropna().unique().tolist()
        for x in shape.chart.series:
            for i, point in enumerate(x.points):
                fill_graph = colors[i].rsplit('#')[-1].lower()
                fill_graph = fill_graph + ('0' * (6 - len(fill_graph)))

                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(fill_graph)


def _extract_tbl_properties(table, tbl_style, row_idx, rowtype):
    '''
    Exctracting table properties.
    '''
    # Get Table rows/columns properties.
    row_text_tag = ['./a:p/a:endParaRPr/a:solidFill/a:srgbClr',
                    './a:p/a:endParaRPr/a:solidFill/a:schemeClr']
    txt_clr_map = {
        './a:p/a:endParaRPr/a:solidFill/a:schemeClr': None,
        './a:p/a:endParaRPr/a:solidFill/a:srgbClr': None
    }

    # Text inside Cell properties information.
    info_txt = table.rows[row_idx].cells[0].text_frame.paragraphs[0]
    if not hasattr(info_txt, 'runs'):
        info_txt.add_run()

    text_dict = {}
    info_txt = info_txt.runs[0].font

    for row_text in row_text_tag:
        row_txt_clr = table.rows[row_idx].cells[0]._tc.txBody.xpath(row_text)
        clr_key = '{}_{}'.format(rowtype, row_text.rsplit(':')[-1])
        if len(row_txt_clr) > 0:
            text_dict[clr_key] = row_txt_clr[0].val
        else:
            text_dict[clr_key] = txt_clr_map.get(row_text)

    text_dict['{}{}'.format(rowtype, '_italic')] = info_txt.italic
    text_dict['{}{}'.format(rowtype, '_bold')] = info_txt.bold
    text_dict['{}{}'.format(rowtype, '_txt_size')] = info_txt.size
    text_dict['{}{}'.format(rowtype, '_font_name')] = info_txt.name
    text_dict['{}{}'.format(rowtype, '_underline')] = info_txt.underline

    # Cell properties information.
    # cell_info = ['./a:srgbClr', './a:schemeClr']
    # cell_info_map = {'./a:schemeClr': None, './a:srgbClr': None}

    cell_dict = {}
    # info_cell = table.rows[row_idx].cells[0]._tc.tcPr.solidFill
    # if info_cell is not None:
    #     for cell_prop in cell_info:
    #         cell_clr = info_cell.xpath(cell_prop)
    #         cell_key = '{}_{}'.format(rowtype, cell_prop.rsplit(':')[-1])
    #         if len(cell_clr) > 0:
    #             cell_dict[cell_key] = cell_clr[0].val
    #         else:
    #             cell_dict[cell_key] = cell_info_map.get(cell_prop)
    tbl_style['text'] = text_dict
    tbl_style['cell'] = cell_dict
    return tbl_style


def _get_table_style(shape, rowtype='header'):
    '''
    Function to get Table style for rows and columns.
    '''
    tbl_style = {}
    table = shape.table
    row_idx = 0 if rowtype == 'header' else 1
    row_idx = 1 if len(table.rows) > 1 else 0
    tbl_style = _extract_tbl_properties(table, tbl_style, row_idx, rowtype)
    # Get Table Header style' from here yet to implement.
    # cell_clr = table.rows[0].cells[0]._tc.tcPr.solidFill.schemeClr.val
    # tbl_style['margin_bottom'] = cell_txt_clr.margin_bottom
    return tbl_style


def _set_text_style(run, table_style, rowtype, text):
    '''
    Function to Apply table styles for rows and columns.
    '''
    run.text = '{}'.format(text)
    rows_text = run.font.fill
    rows_text.solid()

    rgb = table_style.get('{}_{}'.format(rowtype, 'srgbClr'))
    theme_color = table_style.get('{}_{}'.format(rowtype, 'schemeClr'))
    if theme_color:
        rows_text.fore_color.theme_color = theme_color
    else:
        if isinstance(rgb, str):
            run.font.color.rgb = RGBColor.from_string(rgb)
        else:
            try:
                run.font.color.rgb = rgb
            except Exception:
                pass
    run.font.name = table_style['{}{}'.format(rowtype, '_font_name')]
    run.font.size = table_style['{}{}'.format(rowtype, '_txt_size')]
    run.font.bold = table_style['{}{}'.format(rowtype, '_bold')]
    run.font.italic = table_style['{}{}'.format(rowtype, '_italic')]
    run.font.underline = table_style['{}{}'.format(rowtype, '_underline')]


def _extend_table(shape, data, total_rows, total_columns):
    '''
    Function to extend table rows and columns if required.
    '''
    avail_rows = len(shape.table.rows)
    # avail_cols = len(shape.table.columns)

    # col_width = shape.table.columns[0].width
    row_height = shape.table.rows[0].height
    # Extending Table Rows if required based on the data
    while avail_rows < total_rows:

        shape.table.rows._tbl.add_tr(row_height)
        avail_rows += 1

    # Extending Table Columns if required based on the data
    # while avail_cols < total_columns:
    #     shape.table._tbl.tblGrid.add_gridCol(col_width)
    #     avail_cols += 1


def _set_cells_style(shape, cell, cell_prop, row_type, gradient):
    '''
    Set cells properties.
    '''
    props = cell.fill
    props.solid()

    if gradient:
        if isinstance(gradient, str):
            clr = gradient.replace('#', '') + ('0' * (6 - len(gradient.replace('#', ''))))
            props.fore_color.rgb = RGBColor.from_string(clr)
        else:
            props.fore_color.rgb = gradient
        return

    theme_color = cell_prop.get('{}_{}'.format(row_type, 'schemeClr'))
    if theme_color:
        props.fore_color.theme_color = theme_color
        # Setting brightness - Not sure when it will get failed
        bright = shape.table.rows[1].cells[0]._tc.tcPr

        if bright.solidFill.schemeClr.lumMod is not None:
            props.fore_color.brightness = bright.solidFill.schemeClr.lumMod.val - 1.0

        elif bright.solidFill.schemeClr.lumOff is not None:
            props.fore_color.brightness = bright.solidFill.schemeClr.lumOff.val

    else:
        rgb = cell_prop.get('{}_{}'.format(row_type, 'srgbClr'))
        if isinstance(rgb, str):
            props.fore_color.rgb = RGBColor.from_string(rgb)
        else:
            props.fore_color.rgb = rgb


def _color_gradient(info, data):
    '''
    Function to calculate color gradient.
    '''
    color_grad = None
    if info.get('gradient') and 'gradient' in info:
        get_grad = info.get('gradient')
        color_grad = _color.RdYlGn
        numeric = data._get_numeric_data()
        min_data = min([numeric[x].min() for x in numeric.columns])
        max_data = max([numeric[x].max() for x in numeric.columns])
        mean_data = (min_data + max_data) / 2.0
        grad_map = {0: min_data, 1: mean_data, 2: max_data}

        if isinstance(get_grad, str):
            color_grad = getattr(_color, get_grad)

        elif isinstance(get_grad, list) and isinstance(get_grad[0], tuple):
            color_grad = get_grad

        elif isinstance(get_grad, list) and isinstance(get_grad[0], str):
            color_grad = [(grad_map[index], clr) for index, clr in enumerate(get_grad)]
    return color_grad


def table(shape, spec, table_data):
    '''
    Update an existing Table shape with data.
    '''
    # Getting color gradient if required.
    # color_grad = _color_gradient(info, table_data)
    if not spec.get('table', {}).get('data'):
        return
    table_data = table_data.get(spec.table.data, pd.DataFrame())
    table_data = json.loads(table_data.to_json(orient='records'))
    total_rows = len(table_data) + 1
    total_columns = len(table_data[0])

    # Extending table if required.
    _extend_table(shape, table_data, total_rows, total_columns)
    # Fetching Table Style for All Cells and texts.
    table = shape.table
    tbl_style = _get_table_style(shape, rowtype='row')
    for row_num, row in enumerate(table.rows):
        cols = len(row.cells._tr.tc_lst)
        # Extending cells in newly added rows.
        while cols < total_columns:
            row.cells._tr.add_tc()
            cols += 1

        row_type = 'header'
        if row_num == 0:
            continue

        row_data = table_data[row_num - 1]
        if row_num > 0:
            row_type = 'row'

        for col_num, cell in enumerate(row.cells):
            # Getting column name
            col_name = table.cell(0, col_num)
            for curr_cell in cell.text_frame.paragraphs:
                if not curr_cell.text.strip():
                    curr_cell.add_run()
                for run in curr_cell.runs:
                    text_prop = tbl_style.get('text', {})
                    row_txt = ''
                    # gradient = False
                    if col_name.text_frame.text in row_data:
                        row_txt = row_data[col_name.text_frame.text]

                    # if color_grad and isinstance(row_txt, (int, float)):
                    #     gradient = color.gradient(row_txt, color_grad)
                    #     # _set_cells_style(shape, cell, cell_prop, row_type, gradient)
                    _set_text_style(run, text_prop, row_type, row_txt)

# Custom Charts Functions below(Sankey, Treemap, Calendarmap).


def make_element():
    """Function to create element structure."""
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }
    a = ElementMaker(namespace=nsmap['a'], nsmap=nsmap)
    p = ElementMaker(namespace=nsmap['p'], nsmap=nsmap)
    r = ElementMaker(namespace=nsmap['r'], nsmap=nsmap)
    return {'nsmap': nsmap, 'a': a, 'p': p, 'r': r}


def fill_color(**kwargs):
    """
    Return a new color object.

    You may use any one of the following ways of specifying colour:

        color(schemeClr='accent2')             # = second theme color
        color(prstClr='black')                 # = #000000
        color(hslClr=[14400000, 100.0, 50.0])  # = #000080
        color(sysClr='windowText')             # = window text color
        color(scrgbClr=(50000, 50000, 50000))  # = #808080
        color(srgbClr='aaccff')                # = #aaccff

    One or more of these modifiers may be specified:

    - alpha    : '10%' indicates 10% opacity
    - alphaMod : '10%' increased alpha by 10% (50% becomes 55%)
    - alphaOff : '10%' increases alpha by 10 points (50% becomes 60%)
    - blue     : '10%' sets the blue component to 10%
    - blueMod  : '10%' increases blue by 10% (50% becomes 55%)
    - blueOff  : '10%' increases blue by 10 points (50% becomes 60%)
    - comp     : True for opposite hue on the color wheel (e.g. red -> cyan)
    - gamma    : True for the sRGB gamma shift of the input color
    - gray     : True for the grayscale version of the color
    - green    : '10%' sets the green component to 10%
    - greenMod : '10%' increases green by 10% (50% becomes 55%)
    - greenOff : '10%' increases green by 10 points (50% becomes 60%)
    - hue      : '14400000' sets the hue component to 14400000
    - hueMod   : '600000' increases hue by 600000 (14400000 becomes 20000000)
    - hueOff   : '10%' increases hue by 10 points (50% becomes 60%)
    - inv      : True for the inverse color. R, G, B are all inverted
    - invGamma : True for the inverse sRGB gamma shift of the input color
    - lum      : '10%' sets the luminance component to 10%
    - lumMod   : '10%' increases luminance by 10% (50% becomes 55%)
    - lumOff   : '10%' increases luminance by 10 points (50% becomes 60%)
    - red      : '10%' sets the red component to 10%
    - redMod   : '10%' increases red by 10% (50% becomes 55%)
    - redOff   : '10%' increases red by 10 points (50% becomes 60%)
    - sat      : '100000' sets the saturation component to 100%
    - satMod   : '10%' increases saturation by 10% (50% becomes 55%)
    - satOff   : '10%' increases saturation by 10 points (50% becomes 60%)
    - shade    : '10%' is 10% of input color, 90% black
    - tint     : '10%' is 10% of input color, 90% white

    Refer
    <http://msdn.microsoft.com/en-in/library/documentformat.openxml.drawing(v=office.14).aspx>
    """
    hslclr = kwargs.get('hslclr')
    sysclr = kwargs.get('sysclr')
    srgbclr = kwargs.get('srgbclr')
    prstclr = kwargs.get('prstclr')
    scrgbclr = kwargs.get('scrgbclr')
    schemeclr = kwargs.get('schemeclr')

    ns = xmlns('a')
    srgbclr = srgbclr.rsplit('#')[-1].lower()
    srgbclr = srgbclr + ('0' * (6 - len(srgbclr)))
    if schemeclr:
        s = '<a:schemeClr %s val="%s"/>' % (ns, schemeclr)
    elif srgbclr:
        s = '<a:srgbClr %s val="%s"/>' % (ns, srgbclr)
    elif prstclr:
        s = '<a:prstClr %s val="%s"/>' % (ns, prstclr)
    elif hslclr:
        s = '<a:hslClr %s hue="%.0f" sat="%.2f%%" lum="%.2f%%"/>' % (
            (ns,) + tuple(hslclr))
    elif sysclr:
        s = '<a:sysClr %s val="%s"/>' % (ns, sysclr)
    elif scrgbclr:
        s = '<a:scrgbClr %s r="%.0f" g="%.0f" b="%.0f"/>' % ((ns,) + tuple(
            scrgbclr))
    color = objectify.fromstring(s)
    return color


def xmlns(*prefixes):
    """XML ns."""
    elem_schema = make_element()
    return ' '.join('xmlns:%s="%s"' % (pre, elem_schema['nsmap'][pre]) for pre in prefixes)


def call(val, g, group, default):
    """Callback."""
    if callable(val):
        return val(g)
    return default


def pixel_to_inch(pixel):
    """Function to convert Pixel to Inches based on OS."""
    linux_width = 72.0
    windows_width = 96.0
    os_name = platform.system().lower().strip()
    if os_name == 'windows':
        return Inches(pixel / windows_width)
    return Inches(pixel / linux_width)


def rect_css(shape, rect_fill, stroke):
    """Function to add text to shape."""
    rect_fill = rect_fill.rsplit('#')[-1].lower()
    rect_fill = rect_fill + ('0' * (6 - len(rect_fill)))
    border_fill = stroke.rsplit('#')[-1].lower()
    border_fill = border_fill + ('0' * (6 - len(border_fill)))
    fill = shape.fill
    fill.solid()
    # Rectangle border color.
    border = shape.line.fill
    border.solid()
    border.fore_color.rgb = RGBColor.from_string(border_fill)
    fill.fore_color.rgb = RGBColor.from_string(rect_fill)


def add_text_to_shape(shape, text, font_size, txt_fill):
    """Function to add text to shape."""
    min_inc = 13000
    pixel_inch = 10000
    font_size = max(font_size * pixel_inch, min_inc)
    if font_size > min_inc:
        txt_fill = txt_fill.rsplit('#')[-1].lower()
        txt_fill = txt_fill + (txt_fill[0] * (6 - len(txt_fill)))
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.add_run()
        for run in paragraph.runs:
            run.text = text
            shape_txt = run.font
            shape_txt = run.font.fill
            shape_txt.solid()
            run.font.size = '{:.0f}'.format(font_size)
            run.font.color.rgb = RGBColor.from_string(txt_fill)


def cust_shape(x, y, w, h, _id):
    """Custom shapes."""
    _cstmshape = '<p:sp ' + xmlns('p', 'a') + '>'
    _cstmshape = _cstmshape + """<p:nvSpPr>
            <p:cNvPr id='%s' name='%s'/>
            <p:cNvSpPr/>
            <p:nvPr/>
          </p:nvSpPr>
          <p:spPr>
            <a:xfrm>
              <a:off x='%s' y='%s'/>
              <a:ext cx='%s' cy='%s'/>
            </a:xfrm>
            <a:custGeom>
              <a:avLst/>
              <a:gdLst/>
              <a:ahLst/>
              <a:cxnLst/>
              <a:rect l='0' t='0' r='0' b='0'/>
            </a:custGeom>
          </p:spPr>
        </p:sp>"""
    shp = _cstmshape % (_id, 'Freeform %d' % _id, x, y, w, h)
    return objectify.fromstring(shp)


def sankey_calc(data, spec):
    """Create sankey data logic."""
    x0 = spec['x0']
    size = spec['size']
    group = spec['group']
    width = spec['width']
    order = spec.get('order')
    default_color = '#ccfccf'
    default_stroke = '#ffffff'
    attrs = spec.get('attrs', {})
    sort = spec.get('sort', False)
    text = eval(spec.get('text')) if spec.get('text') else None
    fill_color = eval(spec.get('color')) if spec.get('color') else None

    g = data.groupby(group)
    frame = pd.DataFrame({
        'size': g[group[0]].count() if size is None else g[size].sum(),
        'seq': 0 if order is None else order(g),
    })
    frame['width'] = frame['size'] / float(frame['size'].sum()) * width
    frame['fill'] = call(fill_color, g, group, default_color)
    frame['text'] = call(text, g, group, '')
    # Add all attrs to the frame as well
    for key, val in iteritems(attrs):
        frame[key] = call(val, g, group, None)
    if 'stroke' not in attrs:
        frame['stroke'] = default_stroke
    # Compute frame['x'] only after sorting
    if order and sort:
        frame.sort_values('seq', inplace=True)
    frame['x'] = x0 + frame['width'].cumsum() - frame['width']
    return frame


def sankey(shape, spec, data):
    """Draw sankey in Treemap."""
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    # Getting parent shapes
    pxl_to_inch = 10000
    default_thickness = 40
    spec = spec['sankey']
    data = data[spec['data']]
    shapes = shape._parent
    y0 = shape.top
    x0 = shape.left

    width = shape.width
    shape_ids = {'shape': 0}
    height = shape.height
    groups = spec['groups']
    thickness = spec.get('thickness', default_thickness) * pxl_to_inch

    h = (height - (thickness * len(groups))) / (len(groups) - 1) + thickness
    frames = {}
    # Sankey Rectangles and texts.
    sankey_conf = {}
    sankey_conf['x0'] = x0
    sankey_conf['size'] = spec.get('size')
    sankey_conf['width'] = width
    sankey_conf['text'] = spec.get('text')
    sankey_conf['color'] = spec.get('color')
    sankey_conf['attrs'] = spec.get('attrs', {})
    sankey_conf['sort'] = spec.get('sort', False)
    stroke = spec.get('stroke', '#ffffff')
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    elem_schema = make_element()
    for ibar, group in enumerate(groups):
        y = y0 + h * ibar
        sankey_conf['group'] = [group]
        df = frames[group] = sankey_calc(data, sankey_conf)
        # Adding rectangle
        for key, row in df.iterrows():
            shp = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, row['x'], y, row['width'], thickness)
            rect_css(shp, row['fill'], stroke)
            txt_fill = _color.contrast(row['fill'])
            add_text_to_shape(shp, row['text'], spec.get('font-size', 18), txt_fill)

    # Sankey Connection Arcs.
    for ibar, (group1, group2) in enumerate(zip(groups[:-1], groups[1:])):
        sankey_conf['group'] = [group1, group2]
        sankey_conf['sort'] = False
        df = sankey_calc(data, sankey_conf)
        pos = collections.defaultdict(float)
        for key1, row1 in frames[group1].iterrows():
            for key2, row2 in frames[group2].iterrows():
                if (key1, key2) in df.index:
                    row = df.ix[(key1, key2)]
                    y1, y2 = y0 + h * ibar + thickness, y0 + h * (ibar + 1)
                    ym = (y1 + y2) / 2
                    x1 = row1['x'] + pos[0, key1]
                    x2 = row2['x'] + pos[1, key2]

                    _id = shape_ids['shape'] = shape_ids['shape'] + 1
                    shp = cust_shape(0, 0, '{:.0f}'.format(row['width']), '{:.0f}'.format(ym), _id)
                    path = elem_schema['a'].path(
                        w='{:.0f}'.format(row['width']), h='{:.0f}'.format(ym))
                    shp.find('.//a:custGeom', namespaces=elem_schema['nsmap']).append(
                        elem_schema['a'].pathLst(path))
                    path.append(
                        elem_schema['a'].moveTo(elem_schema['a'].pt(
                            x='{:.0f}'.format(x1 + row['width']), y='{:.0f}'.format(y1))))

                    path.append(elem_schema['a'].cubicBezTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x1 + row['width']),
                                            y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x2 + row['width']),
                                            y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x2 + row['width']),
                                            y='{:.0f}'.format(y2))))

                    path.append(elem_schema['a'].lnTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(y2))))

                    path.append(elem_schema['a'].cubicBezTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(y1))))

                    path.append(elem_schema['a'].close())
                    shp.spPr.append(elem_schema['a'].solidFill(fill_color(srgbclr=row['fill'])))
                    shapes._spTree.append(shp)
                    pos[0, key1] += row['width']
                    pos[1, key2] += row['width']


def squarified(x, y, w, h, data):
    """

    Draw a squarified treemap.

    See <http://citeseerx.ist.psu.edu/viewdoc/summary?doi=10.1.1.36.6685>
    Returns a numpy array with (x, y, w, h) for each item in data.

    Examples
    --------
    The result is a 2x2 numpy array::

        >>> squarified(x=0, y=0, w=6, h=4, data=[6, 6, 4, 3, 2, 2, 1])
        array([[ 0.        ,  0.        ,  3.        ,  2.        ],
               [ 0.        ,  2.        ,  3.        ,  2.        ],
               [ 3.        ,  0.        ,  1.71428571,  2.33333333],
               [ 4.71428571,  0.        ,  1.28571429,  2.33333333],
               [ 3.        ,  2.33333333,  1.2       ,  1.66666667],
               [ 4.2       ,  2.33333333,  1.2       ,  1.66666667],
               [ 5.4       ,  2.33333333,  0.6       ,  1.66666667]])

        >>> squarified(x=0, y=0, w=1, h=1, data=[np.nan, 0, 1, 2])
        array([[ 0.        ,  0.        ,  0.        ,  0.        ],
               [ 0.        ,  0.        ,  0.        ,  0.        ],
               [ 0.        ,  0.        ,  0.33333333,  1.        ],
               [ 0.33333333,  0.        ,  0.66666667,  1.        ]])
    """
    w, h = float(w), float(h)
    size = np.nan_to_num(np.array(data).astype(float))
    start, end = 0, len(size)
    result = np.zeros([end, 4])
    if w <= 0 or h <= 0:
        return result

    cumsize = np.insert(size.cumsum(), 0, 0)
    while start < end:
        # We lay out out blocks of rects on either the left or the top edge of
        # the remaining rectangle. But how many rects in the block? We take as
        # many as we can as long as the worst aspect ratio of the block's
        # rectangles keeps improving.

        # This section is (and should be) be heavily optimised. Each operation
        # is run on every element in data.
        last_aspect, newstart = np.Inf, start + 1
        startsize = cumsize[start]
        blockmin = blockmax = size[newstart - 1]
        blocksum = cumsize[newstart] - startsize
        datasum = cumsize[end] - startsize
        ratio = datasum * (h / w if w > h else w / h)
        while True:
            f = blocksum * blocksum / ratio
            aspect = blockmax / f if blockmax > f else f / blockmax
            aspect2 = blockmin / f if blockmin > f else f / blockmin
            if aspect2 > aspect:
                aspect = aspect2
            if aspect <= last_aspect:
                if newstart < end:
                    last_aspect = aspect
                    newstart += 1
                    val = size[newstart - 1]
                    if val < blockmin:
                        blockmin = val
                    if val > blockmax:
                        blockmax = val
                    blocksum += val
                else:
                    break
            else:
                if newstart > start + 1:
                    newstart = newstart - 1
                break

        # Now, lay out the block = start:newstart on the left or top edge.
        block = slice(start, newstart)
        blocksum = cumsize[newstart] - startsize
        scale = blocksum / datasum
        blockcumsize = cumsize[block] - startsize

        if w > h:
            # Layout left-edge, downwards
            r = h / blocksum
            result[block, 0] = x
            result[block, 1] = y + r * blockcumsize
            result[block, 2] = dx = w * scale
            result[block, 3] = r * size[block]
            x, w = x + dx, w - dx
        else:
            # Layout top-edge, rightwards
            r = w / blocksum
            result[block, 0] = x + r * blockcumsize
            result[block, 1] = y
            result[block, 2] = r * size[block]
            result[block, 3] = dy = h * scale
            y, h = y + dy, h - dy

        start = newstart

    return np.nan_to_num(result)


class SubTreemap(object):
    """
    Yield a hierarchical treemap at multiple levels.

    Usage:
        SubTreemap(
            data=data,
            keys=['Parent', 'Child'],
            values={'Value':sum},
            size=lambda x: x['Value'],
            sort=None,
            padding=0,
            aspect=1)

    yields:
        x, y, w, h, (level, data)
    """

    def __init__(self, **args):
        """Default Constructor."""
        self.args = args

    def draw(self, width, height, x=0, y=0, filter={}, level=0):
        """Function to draw rectanfles."""
        # We recursively into each column in `keys` and stop there
        if level >= len(self.args['keys']):
            return

        # Start with the base dataset. Filter by each key applied so far
        summary = self.args['data']
        for key in filter:
            summary = summary[summary[key] == filter[key]]

        # Aggregate by the key up to the current level
        summary = summary.groupby(
            self.args['keys'][:level + 1]
        ).agg(self.args.get('values', {}))
        for key in self.args['keys'][:level + 1]:
            if hasattr(summary, 'reset_index'):
                # Just pop the key out. .reset_index(key) should do this.
                # But on Pandas 0.20.1, this fails
                summary = summary.reset_index([summary.index.names.index(key)])
            else:
                summary[key] = summary.index

        # If specified, sort the aggregated data
        if 'sort' in self.args and self.args['sort'] is callable:
            summary = eval(self.args['sort'])(summary)

        pad = self.args.get('padding', 0)
        aspect = self.args.get('aspect', 1)

        # Find the positions of each box at this level
        key = self.args['keys'][level]
        rows = (summary.to_records() if hasattr(summary, 'to_records') else
                summary)
        rects = squarified(x, y * aspect, width, height * aspect, eval(self.args['size'])(rows))
        for i2, (x2, y2, w2, h2) in enumerate(rects):
            v2 = rows[i2]
            y2, h2 = y2 / aspect, h2 / aspect
            # Ignore invalid boxes generated by Squarified
            if (
                np.isnan([x2, y2, w2, h2]).any() or
                np.isinf([x2, y2, w2, h2]).any() or
                w2 < 0 or h2 < 0
            ):
                continue

            # For each box, dive into the next level
            filter2 = dict(filter)
            filter2.update({key: v2[key]})
            for output in self.draw(w2 - 2 * pad, h2 - 2 * pad, x=x2 + pad, y=y2 + pad,
                                    filter=filter2, level=level + 1):
                yield output

            # Once we've finished yielding smaller boxes, yield the parent box
            yield x2, y2, w2, h2, (level, v2)


def treemap(shape, spec, data):
    """Function to download data as ppt."""
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    shapes = shape._parent
    x0 = shape.left
    y0 = shape.top
    width = shape.width
    height = shape.height
    spec = spec['treemap']
    stroke = spec.get('stroke', '#ffffff')
    spec['data'] = data[spec['data']]
    # Getting rectangle's width and height using `squarified` algorithm.
    treemap_data = SubTreemap(**spec)
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    font_aspect = 14.5
    pixel_inch = 10000
    for x, y, w, h, (l, v) in treemap_data.draw(width, height):
        if l == 0:
            shp = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + x0, y + y0, w, h)
            rect_color = '#cccccc'
            if spec.get('color'):
                rect_color = eval(spec['color'])(v)
            if spec.get('text'):
                text = eval(spec['text'])(v)
            else:
                text = '{}'.format(v[1])
            rect_css(shp, rect_color, stroke)
            txt_fill = _color.contrast(rect_color)
            font_size = min(h, w * font_aspect / fontwidth.fontwidth('{}'.format(text)), pd.np.Inf)
            # Adding text inside rectangles
            add_text_to_shape(shp, text, font_size / pixel_inch, txt_fill)


def calendarmap(shape, spec, data):
    """Draw calendar map in PPT."""
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()

    shapes = shape._parent
    spec = spec['calendarmap']
    data = eval('data.{}'.format(spec['data']))
    pixel_inch = 10000
    size = spec.get('size', None)
    startdate = eval(spec['startdate'])

    label_top = spec.get('label_top', 0) * pixel_inch
    label_left = spec.get('label_left', 0) * pixel_inch

    width = spec['width'] * pixel_inch
    shape_top = label_top + shape.top
    shape_left = label_left + shape.left
    y0 = width + shape_top
    x0 = width + shape_left

    # Deleting the shape
    shape.element.delete()
    # Style
    default_color = '#ffffff'
    default_line_color = '#787C74'
    default_txt_color = '#000000'
    font_size = spec.get('font-size', 12)
    stroke = spec.get('stroke', '#ffffff')
    fill_rect = spec.get('fill', '#cccccc')
    text_color = spec.get('text-color', '#000000')
    # Treat infinities as nans when calculating scale
    scaledata = pd.Series(data).replace([pd.np.inf, -pd.np.inf], pd.np.nan)
    lo_data = spec.get('lo', scaledata.min())
    range_data = spec.get('hi', scaledata.max()) - lo_data
    gradient = spec.get('gradient', _color.RdYlGn)
    color = spec.get('color', lambda v: _color.gradient(
        (float(v) - lo_data) / range_data, gradient) if not pd.isnull(v) else default_color)

    startweekday = (startdate.weekday() - spec.get('weekstart', 0)) % 7
    # Weekday Mean and format
    weekday_mean = pd.Series(
        [scaledata[(x - startweekday) % 7::7].mean() for x in range(7)])
    weekday_format = spec.get('format', '{:,.%df}' % stats.decimals(weekday_mean.values))
    # Weekly Mean and format
    weekly_mean = pd.Series([scaledata[max(0, x):x + 7].mean()
                             for x in range(-startweekday, len(scaledata) + 7, 7)])
    weekly_format = spec.get('format', '{:,.%df}' % stats.decimals(weekly_mean.values))

    # Scale sizes as square roots from 0 to max (not lowest to max -- these
    # should be an absolute scale)
    sizes = width * stats.scale(
        [v ** .5 for v in size], lo=0) if size is not None else [width] * len(scaledata)

    for i, val in enumerate(data):
        nx = (i + startweekday) // 7
        ny = (i + startweekday) % 7
        d = startdate + datetime.timedelta(days=i)
        fill = '#ffffff'
        if not pd.isnull(val):
            fill = color(val)

        shp = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x0 + (width * nx) + (width - sizes[i]) / 2,
            y0 + (width * ny) + (width - sizes[i]) / 2,
            sizes[i], sizes[i])
        rect_css(shp, fill, stroke)
        add_text_to_shape(shp, '%02d' % d.day, font_size, _color.contrast(fill))

        # Draw the boundary lines between months
        if i >= 7 and d.day == 1 and ny > 0:
            vertical_line = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x0 + width * nx, y0 + (width * ny), width, 2 * pixel_inch)
            rect_css(vertical_line, default_line_color, default_line_color)

        if i >= 7 and d.day <= 7 and nx > 0:
            horizontal_line = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x0 + (width * nx), y0 + (width * ny), 2 * pixel_inch, width)
            rect_css(horizontal_line, default_line_color, default_line_color)

        # Adding weekdays text to the chart (left side)
        if i < 7:
            txt = shapes.add_textbox(
                x0 - (width / 2), y0 + (width * ny) + (width / 2), width, width)
            add_text_to_shape(txt, d.strftime('%a')[0], font_size, default_txt_color)

        # Adding months text to the chart (top)
        if d.day <= 7 and ny == 0:
            txt = shapes.add_textbox(
                x0 + (width * nx), y0 - (width / 2), width, width)
            add_text_to_shape(txt, d.strftime('%b %Y'), font_size, default_txt_color)

    if label_top:
        lo_weekly = spec.get('lo', weekly_mean.min())
        range_weekly = spec.get('hi', weekly_mean.max()) - lo_weekly
        for nx, val in enumerate(weekly_mean):
            if not pd.isnull(val):
                w = label_top * ((val - lo_weekly) / range_weekly)
                px = x0 + (width * nx)
                top_bar = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, px, shape_top - w, width, w)
                rect_css(top_bar, fill_rect, stroke)
                top_txt = shapes.add_textbox(px, shape_top - width, width, width)
                add_text_to_shape(
                    top_txt, weekly_format.format(weekly_mean[nx]), font_size, text_color)

    if label_left:
        lo_weekday = spec.get('lo', weekday_mean.min())
        range_weekday = spec.get('hi', weekday_mean.max()) - lo_weekday
        for ny, val in enumerate(weekday_mean):
            if not pd.isnull(val):
                w = label_left * ((val - lo_weekday) / range_weekday)
                bar = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, shape_left - w, y0 + (width * ny), w, width)
                rect_css(bar, fill_rect, stroke)
                left_txt = shapes.add_textbox(shape_left - width, y0 + (width * ny), w, width)
                add_text_to_shape(
                    left_txt, weekday_format.format(weekday_mean[ny]), font_size, text_color)


cmdlist = {
    'text': text,
    'oval': oval,
    'image': image,
    'chart': chart,
    'table': table,
    'sankey': sankey,
    'treemap': treemap,
    'replace': replace,
    'calendarmap': calendarmap
}
